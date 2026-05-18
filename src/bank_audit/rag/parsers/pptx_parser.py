"""PPTX parser: текст со слайдов + speaker notes."""
from __future__ import annotations
import io, re
from pptx import Presentation

from .base import ParsedDoc


def parse_pptx(content: bytes, url: str = "") -> ParsedDoc:
    prs = Presentation(io.BytesIO(content))
    out_lines: list[str] = []
    title: str | None = None

    for i, slide in enumerate(prs.slides, start=1):
        slide_title = None
        slide_lines: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = (shape.text_frame.text or "").strip()
            if not text:
                continue
            if shape == slide.shapes.title or (
                slide.shapes.title is not None and shape == slide.shapes.title
            ):
                slide_title = text
                if title is None:
                    title = text
                continue
            slide_lines.append(text)

        # Speaker notes
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()

        # Заголовок слайда → markdown ##
        out_lines.append(f"\n## Слайд {i}: {slide_title or '(без заголовка)'}\n")
        for line in slide_lines:
            out_lines.append(line.strip())
        if notes:
            out_lines.append(f"\n_Заметки докладчика:_ {notes}\n")

    body = "\n\n".join(out_lines)
    body = re.sub(r"\n{3,}", "\n\n", body).strip()

    return ParsedDoc(
        title=title, text=body, doc_type="pptx",
        meta={"url": url, "slide_count": len(prs.slides),
              "char_count": len(body)},
    )
